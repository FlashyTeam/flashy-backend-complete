from flask import Flask, request, render_template
from pdf2image import convert_from_path
import PyPDF2
import os
from docx import Document
from pptx import Presentation

app = Flask(__name__)

# Ensure the static folder exists for saving images
if not os.path.exists('static'):
    os.makedirs('static')

def chunk_text_by_lines(paragraphs, lines_per_page=25):
    chunks = []
    current_chunk = []
    current_lines = 0
    for p in paragraphs:
        lines_in_p = p.count('\n') + 1
        if current_lines + lines_in_p > lines_per_page:
            chunks.append(current_chunk)
            current_chunk = []
            current_lines = 0
        current_chunk.append(p)
        current_lines += lines_in_p
    if current_chunk:
        chunks.append(current_chunk)
    return chunks

@app.route('/upload', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']
        course = request.form.get('text', 'general').title()
        print(course)
        
        if file:
            file_ext = file.filename.split('.')[-1].lower()
            if file_ext == 'pdf':
                file.save('uploaded_doc.pdf')
                images = convert_from_path('uploaded_doc.pdf',500,poppler_path=r'C:\Program Files\poppler-24.07.0\Library\bin')
                image_paths = []
                for i, image in enumerate(images):
                    image_path = f'static/page_{i + 1}.png'
                    image.save(image_path, 'PNG')
                    image_paths.append(image_path)
                return render_template('preview_pdf.html', images=image_paths, file_ext=file_ext)
            elif file_ext == 'docx':
                file.save('uploaded_doc.docx')
                doc = Document('uploaded_doc.docx')
                paragraphs = [p.text for p in doc.paragraphs]
                pages = chunk_text_by_lines(paragraphs)
                return render_template('preview_text.html', pages=pages, file_ext=file_ext)
            elif file_ext == 'pptx':
                file.save('uploaded_presentation.pptx')
                ppt = Presentation('uploaded_presentation.pptx')
                slides = [slide.shapes.title.text if slide.shapes.title else '' for slide in ppt.slides]
                return render_template('preview_text.html', pages=[slides], file_ext=file_ext)
    return render_template('upload.html')

@app.route('/process', methods=['POST'])
def process_pages():
    selected_pages = request.form.getlist('pages')
    file_ext = request.form.get('file_ext')
    extracted_text = ''

    if file_ext == 'pdf':
        reader = PyPDF2.PdfReader('uploaded_doc.pdf')
        for page_num in selected_pages:
            page = reader.pages[int(page_num) - 1]
            extracted_text += page.extract_text()
        # Delete images after processing
        for i in range(len(reader.pages)):
            image_path = f'static/page_{i + 1}.png'
            if os.path.exists(image_path):
                os.remove(image_path)
        if os.path.exists('uploaded_doc.pdf'):
            os.remove('uploaded_doc.pdf')

    elif file_ext == 'docx':
        doc = Document('uploaded_doc.docx')
        paragraphs = [p.text for p in doc.paragraphs]
        pages = chunk_text_by_lines(paragraphs)
        for page_num in selected_pages:
            extracted_text += '\n'.join(pages[int(page_num) - 1])
        if os.path.exists('uploaded_doc.docx'):
            os.remove('uploaded_doc.docx')

    elif file_ext == 'pptx':
        ppt = Presentation('uploaded_presentation.pptx')
        slides = [slide.shapes.title.text if slide.shapes.title else '' for slide in ppt.slides]
        for page_num in selected_pages:
            extracted_text += slides[int(page_num) - 1] + '\n'
        if os.path.exists('uploaded_presentation.pptx'):
            os.remove('uploaded_presentation.pptx')

    return render_template('result.html', text=extracted_text)

if __name__ == '__main__':
    app.run(debug=True)
