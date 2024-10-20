from flask import Flask, request, jsonify, render_template, redirect, url_for, flash, session
from werkzeug.utils import secure_filename
from flask_bcrypt import Bcrypt
import stripe
import config
from flask_cors import CORS
from pdf2image import convert_from_path
import PyPDF2
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import os
import re
from dotenv import load_dotenv
import google.generativeai as genai
from pymongo import MongoClient
import json
import time

load_dotenv()


app = Flask(__name__)

# Advanced CORS configuration
cors = CORS(app, resources={r"/*": {"origins": "*"}})

app.secret_key = os.getenv('SECRET_KEY')
bcrypt = Bcrypt(app)

# Set your Gemini API key
api_key = os.getenv('API_KEY')
genai.configure(api_key=api_key)

app.config['UPLOAD_FOLDER'] = 'uploads'

# Set up MongoDB
client = MongoClient(os.getenv('MONGO_URI'))
db = client.flashy
flashcards_collection = db.flashcards
users = db.users


# Static folder for storing generated images from pdf
if not os.path.exists('static'):
    os.makedirs('static')


def chunk_text_by_lines(paragraphs, lines_per_page=25):
    chunks = []
    current_chunk = []
    current_lines = 0
    for p in paragraphs:
        lines_in_p = p.count('\n') + 1
        if current_lines + lines_in_p > lines_per_page:
            chunks.append(current_chunk)
            current_chunk = []
            current_lines = 0
        current_chunk.append(p)
        current_lines += lines_in_p
    if current_chunk:
        chunks.append(current_chunk)
    return chunks



if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx'}

def allowed_file(filename):
    return '.' in filename and \
            filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return "Welcome to the Flashy API"

# Route for generated cards for free version
@app.route('/generated_cards_free')
def generated_cards_free():
    return render_template('generated_cards_free.html')

# Route for generated cards for premium version
@app.route('/generated_cards')
def generated_cards():
    return render_template('generated_cards.html')

# Route for viewing flashcard page
@app.route('/view_flashcards')
def view_flashcards():
    return render_template('view_flashcards.html')

# Upload for free version
@app.route('/upload_free_page')
def upload_free_page():
    return render_template('upload_free.html')

# Upload for premium users
@app.route('/upload_page')
def upload_page():
    if 'user_id' not in session:
        flash("Please log in to access your dashboard.")
        return redirect(url_for('login'))
    return render_template('upload.html', username=session['username'], plan=session['plan'])

# Route for fetching courses in request.jsonbase
@app.route('/courses', methods=['POST'])
def get_courses():
    username = request.json.get('username')

    courses = flashcards_collection.distinct('course', {'username': username})
    return jsonify({'courses': courses})


# Route for retrieving flashcards from request.jsonbase
@app.route('/flashcards', methods=['POST'])
def get_flashcards():
    username = request.json.get('username')
    course = request.json.get('selectedCourse')
    
    # Find the document based on username and course
    flashcards_data = flashcards_collection.find_one({'username': username, 'course': course})
    
    if flashcards_data:
        return jsonify({
            'flashcards': flashcards_data.get('flashcards', []),
            'quizzes': flashcards_data.get('quizzes', [])
        })
    return jsonify({'flashcards': [], 'quizzes': []})

# Route for fetching user details
@app.route('/userAccount', methods=['POST'])
def get_userAccount():
    username = request.json.get('username')
    
    # Find the user in the database
    user = users.find_one({'username': username})
    
    if user:
        email = user.get('email')
        plan = user.get('plan')
        return jsonify({'email': email, 'plan': plan}), 200
    else:
        return jsonify({'error': 'User not found'}), 404


# Register user
@app.route('/register', methods=['POST'])
def register():
    username = request.json.get('username')
    email = request.json.get('email')
    password = request.json.get('password')

    hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')

    user_data = {
        "username": username,
        "email": email,
        "password": hashed_password,
        "plan": "premium"
    }

    if users.find_one({"email": email}):
        return jsonify({'error': 'Email already registered!'}), 400

    if users.find_one({"username": username}):
        return jsonify({'error': 'Username already registered!'}), 400

    users.insert_one(user_data)
    
    return jsonify({'message': 'Registration successful!'}), 201


# Login User
@app.route('/login', methods=['POST'])
def login():
    email = request.json.get('email')
    password = request.json.get('password')

    user = users.find_one({"email": email})
    if user and bcrypt.check_password_hash(user['password'], password):
        session['user_id'] = str(user['_id'])
        session['username'] = user['username']
        session['plan'] = user['plan']
        return jsonify({'username': user['username']}), 200
    else:
        return jsonify({'error': 'Invalid credentials!'}), 401


# Logout User
@app.route('/logout')
def logout():
    session.clear()
    return jsonify({'message': 'You have been logged out.'}), 200


# Upload for free version
@app.route('/preview_free_upload', methods=['POST'])
def preview_free_upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in the request'}), 400

    file = request.files['file']
    print('file uploaded')

    if file:
        file_ext = file.filename.split('.')[-1].lower()
        if file_ext == 'pdf':
            file.save('uploaded_doc.pdf')
            images = convert_from_path('uploaded_doc.pdf', 500, poppler_path=r'C:\Program Files\poppler-24.07.0\Library\bin')
            image_paths = []
            for i, image in enumerate(images):
                image_path = f'static/page_{i + 1}.png'
                image.save(image_path, 'PNG')
                image_paths.append(url_for('static', filename=f'page_{i + 1}.png', _external=True))
            return jsonify({'file_ext': 'pdf', 'images': image_paths})
        elif file_ext == 'docx':
            file.save('uploaded_doc.docx')
            doc = Document('uploaded_doc.docx')
            paragraphs = [p.text for p in doc.paragraphs]
            pages = chunk_text_by_lines(paragraphs)
            return jsonify({'file_ext': 'docx', 'pages': pages})
        elif file_ext == 'pptx':
            file.save('uploaded_presentation.pptx')
            ppt = Presentation('uploaded_presentation.pptx')
            slides_content = []
            for slide in ppt.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slides_content.append('\n'.join(slide_text))
            pages = chunk_text_by_lines(slides_content)
            return jsonify({'file_ext': 'pptx', 'pages': pages})

    return jsonify({'error': 'Unsupported file type'}), 400



# Upload for premium users
@app.route('/preview_upload', methods=['POST'])
def preview_upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    print('file uploaded')
    username = request.form.get('username')
    course = request.form.get('course', 'general').title()
    print(course)
    print(username)

    if file:
        file_ext = file.filename.split('.')[-1].lower()
        if file_ext == 'pdf':
            file.save('uploaded_doc.pdf')
            images = convert_from_path('uploaded_doc.pdf', 500, poppler_path=r'C:\Program Files\poppler-24.07.0\Library\bin')
            image_paths = []
            for i, image in enumerate(images):
                image_path = f'static/page_{i + 1}.png'
                image.save(image_path, 'PNG')
                image_paths.append(url_for('static', filename=f'page_{i + 1}.png', _external=True))
            return jsonify({'file_ext': file_ext, 'images': image_paths, 'course': course, 'username': username})
        elif file_ext == 'docx':
            file.save('uploaded_doc.docx')
            doc = Document('uploaded_doc.docx')
            paragraphs = [p.text for p in doc.paragraphs]
            pages = chunk_text_by_lines(paragraphs)
            return jsonify({'file_ext': file_ext, 'pages': pages, 'course': course, 'username': username})
        elif file_ext == 'pptx':
            file.save('uploaded_presentation.pptx')
            ppt = Presentation('uploaded_presentation.pptx')
            slides_content = []
            for slide in ppt.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slides_content.append('\n'.join(slide_text))
            pages = chunk_text_by_lines(slides_content)
            return jsonify({'file_ext': file_ext, 'pages': pages, 'course': course, 'username': username})
        else:
            return jsonify({'error': 'Unsupported file type'}), 400

    return jsonify({'error': 'No file uploaded'}), 400


# Generate cards for free users
@app.route('/process_free', methods=['POST'])
def process_free_pages():
    selected_pages = request.json.get('pages')
    file_ext = request.json.get('file_ext')
    flashcard_number = request.json.get('flashcard_number')

    try:
        extracted_text = ''
        flashcards = ''
        
        if file_ext == 'pdf':
            reader = PyPDF2.PdfReader('uploaded_doc.pdf')
            for page_num in selected_pages:
                page = reader.pages[int(page_num) - 1]
                extracted_text += page.extract_text()
            # Delete images after processing
            for i in range(len(reader.pages)):
                image_path = f'static/page_{i + 1}.png'
                if os.path.exists(image_path):
                    os.remove(image_path)
            if os.path.exists('uploaded_doc.pdf'):
                os.remove('uploaded_doc.pdf')
            print(extracted_text)

        elif file_ext == 'docx':
            doc = Document('uploaded_doc.docx')
            paragraphs = [p.text for p in doc.paragraphs]
            pages = chunk_text_by_lines(paragraphs)
            for page_num in selected_pages:
                extracted_text += '\n'.join(pages[int(page_num) - 1])
            if os.path.exists('uploaded_doc.docx'):
                os.remove('uploaded_doc.docx')
            print(extracted_text)

        elif file_ext == 'pptx':
            ppt = Presentation('uploaded_presentation.pptx')
            slides_content = []
            for slide in ppt.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slides_content.append('\n'.join(slide_text))
            pages = chunk_text_by_lines(slides_content)
            for page_num in selected_pages:
                extracted_text += '\n'.join(pages[int(page_num) - 1])
            if os.path.exists('uploaded_presentation.pptx'):
                os.remove('uploaded_presentation.pptx')

        # Generate flashcards
        flashcards = generate_flashcards(extracted_text, flashcard_number)
        print(flashcards)
        return jsonify({'flashcards': flashcards})

    except Exception as e:
        print(f"Error processing file: {e}")
        return jsonify({'error': 'An error occurred while processing the file.'}), 500



# Generating cards for premium users
@app.route('/process', methods=['POST'])
def process_pages():
    selected_pages = request.json.get('pages')
    flashcard_number = request.json.get('flashcard_number')
    quiz = request.json.get('quiz')
    file_ext = request.json.get('file_ext')
    course = request.json.get('course', 'General')
    username = request.json.get('username')
    print(f"Selected Pages: {selected_pages}")
    print(f"Flashcard Number: {flashcard_number}")
    print(f"Quiz: {quiz}")
    print(f"File Extension: {file_ext}")
    print(f"Course: {course}")
    print(f"Username: {username}")


    try:
        extracted_text = ''
        flashcards = []

        # Process the file according to the extension
        if file_ext == 'pdf':
            reader = PyPDF2.PdfReader('uploaded_doc.pdf')
            for page_num in selected_pages:
                page = reader.pages[int(page_num) - 1]
                extracted_text += page.extract_text()
            # Delete images after processing
            for i in range(len(reader.pages)):
                image_path = f'static/page_{i + 1}.png'
                if os.path.exists(image_path):
                    os.remove(image_path)
            if os.path.exists('uploaded_doc.pdf'):
                os.remove('uploaded_doc.pdf')

        elif file_ext == 'docx':
            doc = Document('uploaded_doc.docx')
            paragraphs = [p.text for p in doc.paragraphs]
            pages = chunk_text_by_lines(paragraphs)
            for page_num in selected_pages:
                extracted_text += '\n'.join(pages[int(page_num) - 1])
            if os.path.exists('uploaded_doc.docx'):
                os.remove('uploaded_doc.docx')

        elif file_ext == 'pptx':
            ppt = Presentation('uploaded_presentation.pptx')
            slides_content = []
            for slide in ppt.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        slide_text.append(shape.text)
                slides_content.append('\n'.join(slide_text))
            pages = chunk_text_by_lines(slides_content)
            for page_num in selected_pages:
                extracted_text += '\n'.join(pages[int(page_num) - 1])
            if os.path.exists('uploaded_presentation.pptx'):
                os.remove('uploaded_presentation.pptx')

        # Generate flashcards
        flashcards = generate_flashcards(extracted_text, flashcard_number)
        print(flashcards)

        # Store or update the flashcards in MongoDB
        flashcards_collection.update_one(
            {'username': username, 'course': course},
            {'$addToSet': {'flashcards': {'$each': flashcards}}},
            upsert=True
        )

        quiz_cards = []
        if quiz == 'yes':
            quiz_number = 5
            quiz_cards = generate_quiz(extracted_text, quiz_number)
            print(quiz_cards)

            # Store or update the quiz cards in MongoDB
            flashcards_collection.update_one(
                {'username': username, 'course': course},
                {'$addToSet': {'quizzes': {'$each': quiz_cards}}},
                upsert=True
            )

        return jsonify({'flashcards': flashcards, 'quiz_cards':quiz_cards})

    except Exception as e:
        # Handle any errors that occur during processing
        return str(e)


# Clean flashcard response
def clean_text(text):
    # Remove asterisks
    text = text.replace('*', '')
    # Remove HTML line breaks
    text = text.replace('<br>', '\n')
    # Remove any other unwanted symbols (e.g., extra spaces)
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    return text.strip()

# Flashcard generation with gemini
def generate_flashcards(text, num_flashcards, retries=3):
    prompt = f"Create {num_flashcards} flashcards from the following text:\n\n{text}\n\n Return answers in json format as an array of objects with 'front' and 'back' keys. Everything should be inline, no backslash n"

    for attempt in range(retries):
        try:
            model = genai.GenerativeModel('gemini-1.5-pro')
            response = model.generate_content(prompt)
            print(response)
            
            if response and response.candidates:
                flashcards_text = response.candidates[0].content.parts[0].text
                
                # Find the start of the JSON array
                json_start = flashcards_text.find("[")
                if json_start != -1:
                    json_text = flashcards_text[json_start:]
                    try:
                        flashcards = json.loads(json_text)
                        # Clean up the flashcards
                        for card in flashcards:
                            card['front'] = clean_text(card['front'])
                            card['back'] = clean_text(card['back'])
                        return flashcards
                    except json.JSONDecodeError as e:
                        print(f"JSON decode error on attempt {attempt + 1}: {e}")
                else:
                    print(f"No JSON array found in the response on attempt {attempt + 1}.")
            else:
                print(f"Error generating flashcards on attempt {attempt + 1}.")
        
        except Exception as e:
            print(f"Unexpected error on attempt {attempt + 1}: {e}")
        
        time.sleep(1)  # Optional: wait 1 second before retrying

    return "Error generating flashcards after multiple attempts."


# Quiz Generation with gemini
def generate_quiz(text, num_quiz, retries=3):
    prompt = f"Create {num_quiz} multiple-choice questions with answers from the following text:\n\n{text}\n\nReturn answers in JSON format as an array of objects with 'question', 'options' (an array of four possible answers), and 'answer' keys. Everything should be inline, no backslash n and no backslash either."

    for attempt in range(retries):
        try:
            model = genai.GenerativeModel('gemini-1.5-pro')
            response = model.generate_content(prompt)
            # print(response)
            
            if response and response.candidates:
                # Extract the text from the response
                quiz_text = response.candidates[0].content.parts[0].text
                
                # Find the start of the JSON array
                json_start = quiz_text.find("[")
                if json_start != -1:
                    json_text = quiz_text[json_start:]
                    try:
                        quiz = json.loads(json_text)
                        return quiz
                    except json.JSONDecodeError as e:
                        print(f"JSON decode error on attempt {attempt + 1}: {e}")
                else:
                    print(f"No JSON array found in the response on attempt {attempt + 1}.")
            else:
                print(f"Error generating flashcards on attempt {attempt + 1}.")
        
        except Exception as e:
            print(f"Unexpected error on attempt {attempt + 1}: {e}")
        
        time.sleep(1)  # Optional: wait 1 second before retrying

    return "Error generating flashcards after multiple attempts."



if __name__ == '__main__':
    # app.run(debug=True)
    app.run(host='0.0.0.0', port=5000)
