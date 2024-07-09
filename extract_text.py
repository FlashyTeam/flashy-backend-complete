from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import os
from gemini import generate_flashcards

def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return '\n'.join(text)
    except Exception as e:
        return f"Error extracting text from DOCX: {str(e)}"

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        return f"Error extracting text from PPTX: {str(e)}"

def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text.append(page.get_text())
        return '\n'.join(text)
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def main():
    # Ensure these paths are correct
    docx_path = "word.docx"
    pptx_path = "ppt.ppt"
    pdf_path = "pdf.pdf"

    # Check if files exist
    if not os.path.exists(docx_path):
        print(f"DOCX file not found at {docx_path}")
    else:
        docx_text = extract_text_from_docx(docx_path)
        print(f"Text from DOCX:\n{docx_text}\n")

    if not os.path.exists(pptx_path):
        print(f"PPTX file not found at {pptx_path}")
    else:
        pptx_text = extract_text_from_pptx(pptx_path)
        print(f"Text from PPTX:\n{pptx_text}\n")

    if not os.path.exists(pdf_path):
        print(f"PDF file not found at {pdf_path}")
    else:
        pdf_text = extract_text_from_pdf(pdf_path)
        print(f"Text from PDF:\n{pdf_text}\n")
        # generate_flashcards(pdf_text, 8)

if __name__ == "__main__":
    main()
